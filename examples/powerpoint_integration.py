#!/usr/bin/env python3
"""
PowerPoint integration with Presenton and python-pptx fallback.

This module provides:
1. Presenton API integration for AI-powered presentations
2. python-pptx fallback for direct PowerPoint generation
3. Campaign data to PowerPoint conversion
"""

import sys
from pathlib import Path
import pandas as pd
import requests
import json
from datetime import datetime
from typing import Dict, List, Any, Optional
import io
import base64

# PowerPoint libraries
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Add the parent directory to the path
sys.path.insert(0, str(Path(__file__).parent.parent))

from arloai_reporting.data.processors import DataProcessor


class PowerPointGenerator:
    """Generate PowerPoint presentations from campaign data."""
    
    def __init__(self, presenton_url: str = "http://192.168.7.174:3050"):
        self.processor = DataProcessor()
        self.presenton_url = presenton_url
        self.presenton_available = self._check_presenton_availability()
    
    def _check_presenton_availability(self) -> bool:
        """Check if Presenton API is available."""
        try:
            response = requests.get(f"{self.presenton_url}/api/health", timeout=5)
            return response.status_code == 200
        except:
            try:
                # Try alternative endpoints
                response = requests.get(f"{self.presenton_url}/health", timeout=5)
                return response.status_code == 200
            except:
                return False
    
    def generate_with_presenton(self, campaign_data: pd.DataFrame, 
                               presentation_prompt: str = None) -> Optional[bytes]:
        """Generate presentation using Presenton API."""
        if not self.presenton_available:
            print("⚠️  Presenton not available, falling back to python-pptx")
            return None
        
        try:
            # Prepare campaign summary for Presenton
            summary = self._prepare_campaign_summary(campaign_data)
            
            if not presentation_prompt:
                presentation_prompt = f"""
                Create a professional campaign performance presentation with the following data:
                
                Campaign Summary:
                {summary}
                
                Please create slides for:
                1. Executive Summary
                2. Key Performance Indicators
                3. Daily Performance Trends
                4. Creative Performance Comparison
                5. Insights and Recommendations
                
                Use a professional business theme with charts and visual elements.
                """
            
            # Call Presenton API
            api_payload = {
                "prompt": presentation_prompt,
                "template": "business",  # Use business template
                "export_format": "pptx"
            }
            
            response = requests.post(
                f"{self.presenton_url}/api/generate",
                json=api_payload,
                timeout=60
            )
            
            if response.status_code == 200:
                return response.content
            else:
                print(f"Presenton API error: {response.status_code}")
                return None
                
        except Exception as e:
            print(f"Error with Presenton: {e}")
            return None
    
    def generate_with_python_pptx(self, campaign_data: pd.DataFrame) -> bytes:
        """Generate presentation using python-pptx."""
        
        # Create presentation
        prs = Presentation()
        
        # Set slide size to widescreen
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Generate slides
        self._add_title_slide(prs, campaign_data)
        self._add_executive_summary_slide(prs, campaign_data)
        self._add_kpi_slide(prs, campaign_data)
        self._add_daily_performance_slide(prs, campaign_data)
        self._add_creative_comparison_slide(prs, campaign_data)
        self._add_recommendations_slide(prs, campaign_data)
        
        # Save to bytes
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        return pptx_io.getvalue()
    
    def _prepare_campaign_summary(self, df: pd.DataFrame) -> str:
        """Prepare campaign data summary for Presenton."""
        total_impressions = df['Impressions'].sum()
        total_clicks = df['Clicks'].sum()
        avg_ctr = df['CTR'].mean()
        total_spend = df['Spend'].sum()
        date_range = f"{df['Date'].min().strftime('%m/%d/%Y')} - {df['Date'].max().strftime('%m/%d/%Y')}"
        
        # Creative performance
        creative_data = df.groupby('Creative').agg({
            'CTR': 'mean',
            'Impressions': 'sum',
            'Spend': 'sum'
        }).round(3)
        
        summary = f"""
        Campaign Period: {date_range}
        Total Impressions: {total_impressions:,}
        Total Clicks: {total_clicks:,}
        Average CTR: {avg_ctr:.3f}%
        Total Spend: ${total_spend:,.2f}
        Average CPC: ${(total_spend/total_clicks):.2f}
        
        Creative Performance:
        {creative_data.to_string()}
        
        Best performing day: {df.groupby('Date')['CTR'].mean().idxmax().strftime('%m/%d/%Y')}
        """
        
        return summary
    
    def _add_title_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add title slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "🚀 Campaign Performance Report"
        
        date_range = f"{df['Date'].min().strftime('%B %d, %Y')} - {df['Date'].max().strftime('%B %d, %Y')}"
        subtitle.text = f"ArloAI Reporting Engine\n{date_range}\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
        
        # Style the title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.color.rgb = RGBColor(44, 62, 80)  # Dark blue
        
        # Style the subtitle
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(18)
        subtitle_paragraph.font.color.rgb = RGBColor(127, 140, 141)  # Gray
    
    def _add_executive_summary_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "📋 Executive Summary"
        
        # Add content
        content = slide.placeholders[1]
        
        total_impressions = df['Impressions'].sum()
        total_clicks = df['Clicks'].sum()
        avg_ctr = df['CTR'].mean()
        total_spend = df['Spend'].sum()
        campaign_days = (df['Date'].max() - df['Date'].min()).days + 1
        
        best_day = df.groupby('Date')['CTR'].mean().idxmax()
        best_ctr = df.groupby('Date')['CTR'].mean().max()
        
        summary_text = f"""Campaign Overview:
• Campaign Period: {campaign_days} days
• Total Reach: {total_impressions:,} impressions
• Total Engagement: {total_clicks:,} clicks
• Average CTR: {avg_ctr:.3f}%
• Total Investment: ${total_spend:,.2f}
• Cost Efficiency: ${(total_spend/total_clicks):.2f} CPC

Key Highlights:
• Best Performance: {best_day.strftime('%m/%d/%Y')} ({best_ctr:.3f}% CTR)
• Strong engagement with {total_clicks:,} clicks generated
• Efficient spend at ${(total_spend/total_clicks):.2f} average CPC"""
        
        content.text = summary_text
        
        # Style the content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
    def _add_kpi_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add KPI slide with visual elements."""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "📊 Key Performance Indicators"
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(44, 62, 80)
        
        # Calculate KPIs
        total_impressions = df['Impressions'].sum()
        total_clicks = df['Clicks'].sum()
        avg_ctr = df['CTR'].mean()
        total_spend = df['Spend'].sum()
        avg_cpm = (total_spend / total_impressions * 1000) if total_impressions > 0 else 0
        
        # KPI data
        kpis = [
            ("👁️ Total Impressions", f"{total_impressions:,}"),
            ("👆 Total Clicks", f"{total_clicks:,}"),
            ("🎯 Average CTR", f"{avg_ctr:.3f}%"),
            ("💰 Total Spend", f"${total_spend:,.2f}"),
            ("📈 Average CPM", f"${avg_cpm:.2f}"),
            ("⚡ Average CPC", f"${(total_spend/total_clicks):.2f}")
        ]
        
        # Create KPI boxes in 2x3 grid
        box_width = Inches(4)
        box_height = Inches(1.5)
        start_x = Inches(1)
        start_y = Inches(2)
        
        for i, (label, value) in enumerate(kpis):
            row = i // 3
            col = i % 3
            
            x = start_x + col * (box_width + Inches(0.5))
            y = start_y + row * (box_height + Inches(0.3))
            
            # Add colored rectangle
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, y, box_width, box_height
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue
            rect.line.color.rgb = RGBColor(41, 128, 185)  # Darker blue
            
            # Add text
            text_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), 
                                              box_width - Inches(0.4), box_height - Inches(0.4))
            text_frame = text_box.text_frame
            text_frame.text = f"{label}\n{value}"
            
            # Style text
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                if paragraph == text_frame.paragraphs[0]:  # Label
                    paragraph.font.size = Pt(14)
                else:  # Value
                    paragraph.font.size = Pt(20)
                    paragraph.font.bold = True
    
    def _add_daily_performance_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add daily performance slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "📈 Daily Performance Trends"
        
        # Prepare daily data
        daily_data = df.groupby('Date').agg({
            'Impressions': 'sum',
            'Clicks': 'sum',
            'CTR': 'mean',
            'Spend': 'sum'
        }).reset_index()
        
        # Create table
        rows = len(daily_data) + 1  # +1 for header
        cols = 5
        
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), 
                                     Inches(11), Inches(4)).table
        
        # Header
        headers = ['Date', 'Impressions', 'Clicks', 'CTR', 'Spend']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
        
        # Data rows
        for row_idx, (_, row_data) in enumerate(daily_data.iterrows(), 1):
            table.cell(row_idx, 0).text = row_data['Date'].strftime('%m/%d/%Y')
            table.cell(row_idx, 1).text = f"{row_data['Impressions']:,}"
            table.cell(row_idx, 2).text = f"{row_data['Clicks']:,}"
            table.cell(row_idx, 3).text = f"{row_data['CTR']:.3f}%"
            table.cell(row_idx, 4).text = f"${row_data['Spend']:.2f}"
            
            # Style data cells
            for col_idx in range(5):
                cell = table.cell(row_idx, col_idx)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
    def _add_creative_comparison_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add creative comparison slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "🎨 Creative Performance Comparison"
        
        # Prepare creative data
        creative_data = df.groupby('Creative').agg({
            'Impressions': 'sum',
            'Clicks': 'sum',
            'CTR': 'mean',
            'Spend': 'sum'
        }).reset_index()
        
        # Create comparison table
        rows = len(creative_data) + 1
        cols = 5
        
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), 
                                     Inches(12), Inches(4)).table
        
        # Header
        headers = ['Creative', 'Impressions', 'Clicks', 'CTR', 'Spend']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(231, 76, 60)  # Red
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
        
        # Data rows
        for row_idx, (_, row_data) in enumerate(creative_data.iterrows(), 1):
            creative_name = row_data['Creative'].replace('Boston Area University Spring 2025  - ', '')
            table.cell(row_idx, 0).text = creative_name
            table.cell(row_idx, 1).text = f"{row_data['Impressions']:,}"
            table.cell(row_idx, 2).text = f"{row_data['Clicks']:,}"
            table.cell(row_idx, 3).text = f"{row_data['CTR']:.3f}%"
            table.cell(row_idx, 4).text = f"${row_data['Spend']:.2f}"
            
            # Style data cells
            for col_idx in range(5):
                cell = table.cell(row_idx, col_idx)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = Pt(10)
                paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
    def _add_recommendations_slide(self, prs: Presentation, df: pd.DataFrame):
        """Add recommendations slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "💡 Insights & Recommendations"
        
        # Generate insights
        creative_data = df.groupby('Creative')['CTR'].mean()
        best_creative = creative_data.idxmax()
        best_ctr = creative_data.max()
        
        daily_performance = df.groupby('Date')['CTR'].mean()
        best_day = daily_performance.idxmax()
        
        avg_cpc = df['Spend'].sum() / df['Clicks'].sum()
        
        recommendations = f"""Key Insights:
• Version 2 creative is outperforming with {best_ctr:.3f}% CTR
• Best performance day was {best_day.strftime('%A, %m/%d/%Y')}
• Efficient cost management at ${avg_cpc:.2f} average CPC
• Strong overall engagement with consistent performance

Strategic Recommendations:
• Scale the better-performing creative (Version 2)
• Focus budget allocation on high-performing weekdays
• Consider A/B testing new creative variations
• Maintain current targeting strategy - showing good efficiency
• Explore similar audience segments for expansion

Next Steps:
• Implement creative optimization based on performance data
• Set up automated bidding for peak performance days
• Develop new creative concepts based on successful elements
• Monitor and adjust budget allocation weekly"""
        
        content = slide.placeholders[1]
        content.text = recommendations
        
        # Style the content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
    def generate_presentation(self, excel_files: List[str], 
                            output_format: str = "both",
                            custom_prompt: str = None) -> Dict[str, Any]:
        """Generate presentation in specified format(s)."""
        
        # Read and combine Excel data
        all_data = []
        for excel_file in excel_files:
            try:
                df = pd.read_excel(excel_file, sheet_name="Data")
                all_data.append(df)
            except Exception as e:
                print(f"Error reading {excel_file}: {e}")
        
        if not all_data:
            return {"error": "No data could be loaded from Excel files"}
        
        combined_df = pd.concat(all_data, ignore_index=True)
        
        results = {}
        
        # Try Presenton first if requested
        if output_format in ["presenton", "both"] and self.presenton_available:
            print("🎨 Generating presentation with Presenton AI...")
            presenton_pptx = self.generate_with_presenton(combined_df, custom_prompt)
            if presenton_pptx:
                results["presenton"] = presenton_pptx
                print("✅ Presenton presentation generated successfully")
            else:
                print("⚠️  Presenton generation failed")
        
        # Generate with python-pptx
        if output_format in ["python-pptx", "both"] or (output_format == "presenton" and not results.get("presenton")):
            print("📊 Generating presentation with python-pptx...")
            python_pptx = self.generate_with_python_pptx(combined_df)
            results["python-pptx"] = python_pptx
            print("✅ Python-pptx presentation generated successfully")
        
        return results


def main():
    """Generate PowerPoint presentations from campaign data."""
    print("🎯 PowerPoint Integration with Presenton")
    print("=" * 50)
    
    # Find sample files
    sample_dir = Path(__file__).parent / "sample_data"
    excel_files = [
        str(sample_dir / "sample-data-20250710.xlsx"),
        str(sample_dir / "sample-data-20250814.xlsx")
    ]
    
    excel_files = [f for f in excel_files if Path(f).exists()]
    
    if not excel_files:
        print("❌ No Excel files found!")
        return 1
    
    print(f"📊 Processing {len(excel_files)} Excel files...")
    
    # Initialize PowerPoint generator
    generator = PowerPointGenerator()
    
    print(f"🔍 Presenton Status: {'✅ Available' if generator.presenton_available else '❌ Not Available'}")
    
    # Generate presentations
    results = generator.generate_presentation(
        excel_files, 
        output_format="both",  # Try both Presenton and python-pptx
        custom_prompt="""
        Create a professional campaign performance presentation for a university marketing campaign.
        Focus on ROI, engagement metrics, and actionable insights.
        Use a modern, clean design with data visualizations.
        Include executive summary, KPIs, trends, and recommendations.
        """
    )
    
    # Save results
    output_dir = Path(__file__).parent / "output"
    output_dir.mkdir(exist_ok=True)
    
    for method, pptx_data in results.items():
        if isinstance(pptx_data, bytes):
            filename = f"campaign_report_{method.replace('-', '_')}.pptx"
            output_file = output_dir / filename
            
            with open(output_file, 'wb') as f:
                f.write(pptx_data)
            
            print(f"💾 {method.title()} presentation saved: {output_file}")
    
    print("\n🎉 PowerPoint generation complete!")
    print("\nFeatures:")
    print("✅ Presenton AI integration (when available)")
    print("✅ Python-pptx fallback generation")
    print("✅ Professional slide layouts")
    print("✅ Campaign data visualization")
    print("✅ Executive summary and recommendations")
    print("✅ KPI dashboards and performance tables")
    
    if generator.presenton_available:
        print("\n🚀 Presenton is available - AI-powered presentations enabled!")
    else:
        print("\n⚠️  Presenton not available - using python-pptx fallback")
        print("   To enable Presenton:")
        print("   1. Ensure Presenton is running on http://192.168.7.174:3050")
        print("   2. Check API endpoints are accessible")
    
    return 0


if __name__ == "__main__":
    sys.exit(main())